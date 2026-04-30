from http.server import BaseHTTPRequestHandler
import json
import os
import io
import base64
from pptx import Presentation
from pptx.util import Pt
import copy

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")

REPLACEMENTS = {
    "{{nombre_servicio}}": "nombre_servicio",
    "{{descripcion}}":     "descripcion",
    "{{plazo}}":           "plazo",
    "{{incluye}}":         "incluye",
    "{{entregable_1}}":    "entregable_1",
    "{{entregable_2}}":    "entregable_2",
    "{{entregable_3}}":    "entregable_3",
    "{{entregable_4}}":    "entregable_4",
    "{{entregable_5}}":    "entregable_5",
    "{{nota_aclaracion}}": "nota_aclaracion",
    "{{precio}}":          "precio",
    "{{timeline_1}}":      "timeline_1",
    "{{timeline_2}}":      "timeline_2",
    "{{timeline_3}}":      "timeline_3",
}

# These are the actual texts in Gabi's template that map to each field
TEMPLATE_TEXTS = {
    "nombre_servicio": "Creación de Contenido (Video & Foto)+ Brand assesement",
    "descripcion":     "Creación de contenido audiovisual en formato vertical, optimizado para redes sociales, alineado con la estrategia de la marca y pensado para facilitar su ejecución y publicación.",
    "plazo":           "6 semanas",
    "incluye":         "Acompañamiento estratégico para ayudar a una marca a identificar su posición real en el mercado, clarificar su voz y diseñar una hoja de ruta de contenido accionable. A través de un proceso estructurado de reuniones y feedback, definimos dónde está la marca hoy, dónde debería estar, y cómo llegar ahí",
    "entregable_1":    "Brand Discovery Document: análisis de competidores, oportunidad en el mercado digital, tono de voz, do's & don'ts",
    "entregable_2":    "Content Guide (DIY): guía de contenido que el propio equipo de la marca puede producir",
    "entregable_3":    "Filmación y edición de 6 videos (Reels) con contenido previamente definido,",
    "entregable_4":    "60 fotografías de producto y de situación, pensadas para feed, stories y otros usos digitales.",
    "entregable_5":    "Dirección de arte, asegurando coherencia visual, estética de marca y calidad en todo el contenido.",
    "nota_aclaracion": "CON ESTE PACK SE PREVEE CUBRIR EL CONTENIDO PARA UN MÍNIMO DE 3 MESES.",
    "precio":          "2.500€",
    "timeline_1":      "Brand discovery",
    "timeline_2":      "Producción contenido",
    "timeline_3":      "Entrega contenidos",
}

def replace_in_run(run, old, new):
    if old in run.text:
        run.text = run.text.replace(old, new)
        return True
    return False

def replace_in_shape(shape, data):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for field, key in TEMPLATE_TEXTS.items():
                val = data.get(key, "")
                if val and field in run.text:
                    run.text = run.text.replace(field, val)

def replace_text_in_presentation(prs, data):
    """Replace template texts with user data in all slides."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Work at paragraph level to handle split runs
                    full_text = "".join(r.text for r in para.runs)
                    for template_text, field_key in TEMPLATE_TEXTS.items():
                        new_val = data.get(field_key, "")
                        if new_val and template_text in full_text:
                            # Replace in first run that contains part of the text
                            for run in para.runs:
                                if template_text in run.text:
                                    run.text = run.text.replace(template_text, new_val)
                                    break
                            # Handle split across runs - reconstruct
                            full_text = full_text.replace(template_text, new_val)

def generate_pptx(data):
    prs = Presentation(TEMPLATE_PATH)
    replace_text_in_presentation(prs, data)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

class handler(BaseHTTPRequestHandler):
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            body = self.rfile.read(length)
            data = json.loads(body)

            pptx_bytes = generate_pptx(data)
            b64 = base64.b64encode(pptx_bytes).decode("utf-8")

            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            self.wfile.write(json.dumps({"file": b64, "filename": f"Bright_Propuesta_{data.get('nombre_servicio','').replace(' ','_')[:30]}.pptx"}).encode())

        except Exception as e:
            self.send_response(500)
            self.send_header("Content-Type", "application/json")
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(e)}).encode())
